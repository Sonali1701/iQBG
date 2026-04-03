import os
import io
import re
import math
import zipfile
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WHITE_THRESHOLD = 240
CROP_PAD_PX = 20
TARGET_FRACTION = 0.325
ANCHOR_SHAPE_NAME = "SCREENSHOT_BOX"
TEMPLATE_NAME = "Template Format.pptx"

def remove_white(img, thr=WHITE_THRESHOLD):
    img = img.convert('RGBA')
    data = img.getdata()
    new = []
    for (r,g,b,a) in data:
        if r >= thr and g >= thr and b >= thr:
            new.append((255,255,255,0))
        else:
            new.append((r,g,b,a))
    img.putdata(new)
    return img

def crop_to_content_rgba(img_rgba, pad_px=CROP_PAD_PX):
    alpha = img_rgba.split()[-1]
    bbox = alpha.getbbox()
    if not bbox:
        return img_rgba
    x0, y0, x1, y1 = bbox
    x0 = max(0, x0 - pad_px)
    y0 = max(0, y0 - pad_px)
    x1 = min(img_rgba.width,  x1 + pad_px)
    y1 = min(img_rgba.height, y1 + pad_px)
    return img_rgba.crop((x0, y0, x1, y1))

def compute_size_for_area_fraction(slide_w, slide_h, left, top, aspect, target_fraction):
    slide_area = slide_w * slide_h
    desired_area = slide_area * target_fraction

    w = math.sqrt(desired_area * aspect)
    h = math.sqrt(desired_area / aspect)

    max_w = max(1, slide_w - left)
    max_h = max(1, slide_h - top)

    if w > max_w or h > max_h:
        scale = min(max_w / w, max_h / h)
        w *= scale
        h *= scale

    return int(w), int(h)

def iter_all_shapes(shapes):
    for shp in shapes:
        yield shp
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            for inner in iter_all_shapes(shp.shapes):
                yield inner

def find_anchor_on_slide(slide, name):
    for shp in iter_all_shapes(slide.shapes):
        if getattr(shp, "name", "").strip() == name:
            return shp
    return None

def build_single_ppt(image_paths, transparent, output_pptx_path):
    template_path = os.path.join(os.path.dirname(__file__), TEMPLATE_NAME)
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template not found at: {template_path}")
        
    prs = Presentation(template_path)
    if len(prs.slides) == 0:
        raise ValueError("Template has no slides")
        
    template_layout = prs.slides[0].slide_layout

    def ensure_slide(idx):
        while len(prs.slides) <= idx:
            prs.slides.add_slide(template_layout)
        return prs.slides[idx]

    anchor = find_anchor_on_slide(prs.slides[0], ANCHOR_SHAPE_NAME)
    if not anchor:
        raise ValueError(f"Could not find shape named '{ANCHOR_SHAPE_NAME}' on first slide.")
        
    ANCHOR_LEFT = anchor.left
    ANCHOR_TOP = anchor.top
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for i, img_path in enumerate(image_paths):
        slide = ensure_slide(i)
        
        with Image.open(img_path) as im:
            if transparent:
                cleaned = remove_white(im, thr=WHITE_THRESHOLD)
                cleaned = crop_to_content_rgba(cleaned, pad_px=CROP_PAD_PX)
            else:
                cleaned = im.convert('RGBA')

            w_px, h_px = cleaned.size
            if h_px == 0:
                continue
            aspect = w_px / h_px

            w_emu, h_emu = compute_size_for_area_fraction(
                slide_w, slide_h, ANCHOR_LEFT, ANCHOR_TOP, aspect, TARGET_FRACTION
            )

            img_byte_arr = io.BytesIO()
            cleaned.save(img_byte_arr, format='PNG')
            img_byte_arr.seek(0)
            
            slide.shapes.add_picture(img_byte_arr, ANCHOR_LEFT, ANCHOR_TOP, width=w_emu, height=h_emu)

    prs.save(output_pptx_path)

def generate_ppts(job_dir, language_mode, transparent):
    """
    language_mode: "english", "hindi", or "both" (single language mode just passes "english" or "hindi")
    Returns a list of local paths to the generated files.
    """
    def extract_q_num(fname):
        m = re.search(r'_Q(\d+)\.png$', fname)
        return int(m.group(1)) if m else 999999

    eng_images = []
    hin_images = []
    for root, _, files in os.walk(job_dir):
        for f in files:
            if f.endswith(".png"):
                if f.startswith("QUES_ENG"):
                    eng_images.append(os.path.join(root, f))
                elif f.startswith("QUES_HIN"):
                    hin_images.append(os.path.join(root, f))

    # Sort numerically by Q index
    eng_images.sort(key=extract_q_num)
    hin_images.sort(key=extract_q_num)

    generated_files = []

    if language_mode in ["english", "both"]:
        if eng_images:
            out_pptx = os.path.join(job_dir, "English_Questions.pptx")
            build_single_ppt(eng_images, transparent, out_pptx)
            generated_files.append(out_pptx)

    if language_mode in ["hindi", "both"]:
        if hin_images:
            out_pptx = os.path.join(job_dir, "Hindi_Questions.pptx")
            build_single_ppt(hin_images, transparent, out_pptx)
            generated_files.append(out_pptx)

    return generated_files
